import os
import tempfile
import base64
from io import BytesIO
from pptx import Presentation
from utils.pptx_extractors import (
    extract_master_properties,
    extract_shape_properties
)
from utils.style_extractor import (
    extract_placeholder_complete_styling,
    apply_placeholder_styling
)
from utils.slide_layout_extractor import extract_all_slides_as_layouts


class PPTXService:
    def __init__(self):
        self.stored_rules = None
        self.template_path = None
        self.uploaded_images = {}
        self.image_order = []  # track insertion order for image_index lookup
    
    def cleanup_template(self):
        if self.template_path and os.path.exists(self.template_path):
            try:
                os.unlink(self.template_path)
                self.template_path = None
            except Exception as e:
                print(f"error cleaning up template: {e}")
    
    def load_template_file(self, file_storage, layouts, slide_size):
        """load a template file with pre-extracted layouts (used when loading from DB)"""
        self.cleanup_template()
        self.uploaded_images = {}
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            file_storage.save(tmp.name)
            self.template_path = tmp.name
        
        # store the provided rules without re-extracting
        self.stored_rules = {
            'layouts': layouts,
            'slide_size': slide_size,
            'theme_data': None,
            'customTheme': None,
        }
        
        print(f"\ntemplate loaded with {len(layouts)} layouts from database")
    
    def extract_rules_from_file(self, file_storage):
        self.cleanup_template()
        self.uploaded_images = {}
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            file_storage.save(tmp.name)
            self.template_path = tmp.name
        
        try:
            prs = Presentation(self.template_path)
            
            if len(prs.slides) == 0:
                raise ValueError('The uploaded presentation has no slides. Please upload a presentation with at least one example slide.')
            
            extraction_result = extract_all_slides_as_layouts(prs, self.template_path)
            layouts = extraction_result['layouts']
            theme_data = extraction_result.get('theme_data')
            
            print(f"\nextracted {len(layouts)} rigid layout templates")
            
            # STRICT VALIDATION: Fail fast if critical design specs are missing
            self._validate_complete_extraction(prs, layouts, theme_data)
            
            from services.ai_service import AIService
            
            ai_service = AIService()
            categorization_result = ai_service.categorize_layouts(layouts)
            
            rules = {
                'slide_size': {
                    'width': prs.slide_width,
                    'height': prs.slide_height
                },
                'layouts': categorization_result['layouts'],
                'theme_data': theme_data,
                'extraction_method': 'rigid_slide_templates'
            }
            
            text_placeholder_count = sum(
                len([p for p in layout['placeholders'] if p['type'] == 'text'])
                for layout in layouts
            )
            image_placeholder_count = sum(
                len([p for p in layout['placeholders'] if p['type'] == 'image'])
                for layout in layouts
            )
            
            print(f"  total text placeholders: {text_placeholder_count}")
            print(f"  total image placeholders: {image_placeholder_count}")
            
            # Log theme data extraction
            if theme_data:
                if theme_data.get('color_scheme'):
                    print(f"  ✓ theme: {len(theme_data['color_scheme'])} colors extracted")
                if theme_data.get('format_scheme'):
                    print(f"  ✓ theme: format styles extracted")
                if theme_data.get('theme_raw'):
                    print(f"  ✓ theme: raw theme XML preserved")
            
            self.stored_rules = rules
            return rules
            
        except Exception as e:
            self.cleanup_template()
            raise
    
    def get_stored_rules(self):
        if self.stored_rules is None:
            raise ValueError('No rules stored. Please upload a PowerPoint file first.')
        return self.stored_rules
    
    def _validate_complete_extraction(self, prs, layouts, theme_data):
        """
        Strict validation: Ensure all critical design specifications were successfully extracted.
        Fails fast if any required styling information is missing.
        """
        errors = []
        
        # 1. Validate slide dimensions
        if not prs.slide_width or not prs.slide_height:
            errors.append("Failed to extract slide dimensions (width/height)")
        
        # 2. Validate theme data
        if not theme_data:
            errors.append("Failed to extract theme data from presentation")
        else:
            # Validate fonts
            if not theme_data.get('fonts'):
                errors.append("Failed to extract font information from theme")
            else:
                fonts = theme_data['fonts']
                if not fonts.get('title') or not fonts['title'].get('name'):
                    errors.append("Failed to extract title font name from theme")
                if not fonts.get('body') or not fonts['body'].get('name'):
                    errors.append("Failed to extract body font name from theme")
            
            # Validate color scheme
            if not theme_data.get('color_scheme'):
                errors.append("Failed to extract color scheme from theme")
        
        # 3. Validate layouts
        if not layouts or len(layouts) == 0:
            errors.append("No layouts were successfully extracted from presentation")
        
        for i, layout in enumerate(layouts):
            layout_name = layout.get('name', f'Layout {i}')
            
            # Validate layout has placeholders
            if 'placeholders' not in layout:
                errors.append(f"Layout '{layout_name}': Missing placeholders array")
                continue
            
            # Validate each placeholder has complete properties
            for ph in layout['placeholders']:
                ph_idx = ph.get('idx', '?')
                
                # Check for required fields
                if 'type' not in ph:
                    errors.append(f"Layout '{layout_name}', placeholder {ph_idx}: Missing type")
                if 'properties' not in ph:
                    errors.append(f"Layout '{layout_name}', placeholder {ph_idx}: Missing properties")
                    continue
                
                props = ph['properties']
                
                # Validate position data
                if 'position' not in props:
                    errors.append(f"Layout '{layout_name}', placeholder {ph_idx}: Missing position data")
                else:
                    pos = props['position']
                    required_pos_fields = ['left', 'top', 'width', 'height']
                    for field in required_pos_fields:
                        if field not in pos or pos[field] is None:
                            errors.append(f"Layout '{layout_name}', placeholder {ph_idx}: Missing position.{field}")
                
                # For text placeholders, validate font properties
                if ph.get('type') == 'text':
                    if 'font_props' not in props:
                        errors.append(f"Layout '{layout_name}', text placeholder {ph_idx}: Missing font properties")
                    else:
                        font_props = props['font_props']
                        if not font_props.get('name'):
                            errors.append(f"Layout '{layout_name}', text placeholder {ph_idx}: Missing font name")
                        if not font_props.get('size'):
                            errors.append(f"Layout '{layout_name}', text placeholder {ph_idx}: Missing font size")
        
        # If any errors found, fail fast
        if errors:
            error_message = "TEMPLATE EXTRACTION FAILED - Incomplete design specifications:\n\n"
            error_message += "\n".join(f"  • {err}" for err in errors)
            error_message += "\n\nPlease upload a valid PowerPoint template with complete formatting information."
            error_message += "\nEnsure the template has:"
            error_message += "\n  - Defined theme with fonts and colors"
            error_message += "\n  - Properly formatted slide layouts"
            error_message += "\n  - Text placeholders with font specifications"
            raise ValueError(error_message)
        
        print("  all critical design specifications validated successfully")
    
    def update_stored_rules(self, rules):
        """update the stored rules (used for modifying layout properties like is_special)"""
        self.stored_rules = rules
    
    def store_image(self, filename, image_data_base64):
        self.uploaded_images[filename] = image_data_base64
        if filename not in self.image_order:
            self.image_order.append(filename)
    
    def clear_images(self):
        self.uploaded_images = {}
        self.image_order = []
    
    def _apply_theme_overrides(self, custom_theme):
        """
        Apply custom theme overrides to the stored rules.
        This modifies layout templates to use custom fonts and colors.
        """
        if not custom_theme or not self.stored_rules:
            return
        
        custom_fonts = custom_theme.get('fonts', {})
        custom_colors = custom_theme.get('colors', {})
        
        if custom_fonts:
            print(f"  ✓ overriding fonts:")
            if 'title' in custom_fonts:
                print(f"    - title: {custom_fonts['title'].get('name')} {custom_fonts['title'].get('size')}pt")
            if 'body' in custom_fonts:
                print(f"    - body: {custom_fonts['body'].get('name')} {custom_fonts['body'].get('size')}pt")
        
        if custom_colors:
            print(f"  ✓ overriding {len(custom_colors)} theme colors")
        
        # Update all layout templates with custom theme
        layouts = self.stored_rules.get('layouts', [])
        for layout in layouts:
            for placeholder in layout.get('placeholders', []):
                props = placeholder.get('properties', {})
                
                # Override font properties in text placeholders
                if placeholder.get('type') == 'text' and 'font_props' in props:
                    font_props = props['font_props']
                    
                    # Determine if this is a title or body based on name/position
                    is_title = 'title' in placeholder.get('name', '').lower()
                    
                    if is_title and 'title' in custom_fonts:
                        if 'name' in custom_fonts['title']:
                            font_props['name'] = custom_fonts['title']['name']
                        if 'size' in custom_fonts['title']:
                            font_props['size'] = custom_fonts['title']['size']
                    elif not is_title and 'body' in custom_fonts:
                        if 'name' in custom_fonts['body']:
                            font_props['name'] = custom_fonts['body']['name']
                        if 'size' in custom_fonts['body']:
                            font_props['size'] = custom_fonts['body']['size']
                    
                    # Override color if specified and matches theme color names
                    if 'color' in font_props and custom_colors:
                        # For now, apply primary text color (dk1) if available
                        if 'dk1' in custom_colors:
                            color_info = custom_colors['dk1']
                            if color_info.get('type') == 'rgb' and 'value' in color_info:
                                font_props['color'] = {
                                    'type': 'rgb',
                                    'rgb': color_info['value']
                                }
        
        print(f"  theme overrides applied to {len(layouts)} layouts")
    
    def _validate_and_fix_shape_positions(self, presentation):
        """
        validate that all shapes are positioned within slide bounds.
        fix any shapes that extend beyond the slide edges.
        uses deterministic python-pptx calculations.
        """
        # get slide dimensions (in EMUs)
        # standard 16:9 slide: 10" × 5.625" = 9144000 × 5143500 EMUs
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        
        print(f"  slide dimensions: {slide_width / 914400:.2f}\" × {slide_height / 914400:.2f}\" ({slide_width} × {slide_height} EMUs)")
        
        issues_found = 0
        issues_fixed = 0
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            for shape in slide.shapes:
                try:
                    # skip shapes without position (e.g., background)
                    if not hasattr(shape, 'left') or not hasattr(shape, 'top'):
                        continue
                    
                    # get shape bounds
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    right = left + width
                    bottom = top + height
                    
                    # check if shape is out of bounds
                    out_of_bounds = False
                    adjustments = []
                    
                    if left < 0:
                        out_of_bounds = True
                        adjustments.append(f"left {left} → 0")
                        shape.left = 0
                    
                    if top < 0:
                        out_of_bounds = True
                        adjustments.append(f"top {top} → 0")
                        shape.top = 0
                    
                    if right > slide_width:
                        out_of_bounds = True
                        # if shape extends beyond right edge, move it left
                        if width < slide_width:
                            new_left = slide_width - width
                            adjustments.append(f"left {left} → {new_left} (right edge at {right} > {slide_width})")
                            shape.left = new_left
                        else:
                            # shape is too wide for slide - shrink it
                            new_width = slide_width - 100000  # leave small margin
                            adjustments.append(f"width {width} → {new_width} (too wide)")
                            shape.width = new_width
                            shape.left = 50000  # small margin
                    
                    if bottom > slide_height:
                        out_of_bounds = True
                        # if shape extends beyond bottom edge, move it up
                        if height < slide_height:
                            new_top = slide_height - height
                            adjustments.append(f"top {top} → {new_top} (bottom edge at {bottom} > {slide_height})")
                            shape.top = new_top
                        else:
                            # shape is too tall for slide - shrink it
                            new_height = slide_height - 100000  # leave small margin
                            adjustments.append(f"height {height} → {new_height} (too tall)")
                            shape.height = new_height
                            shape.top = 50000  # small margin
                    
                    if out_of_bounds:
                        issues_found += 1
                        issues_fixed += 1
                        shape_name = getattr(shape, 'name', 'Unknown')
                        print(f"    slide {slide_num}, shape '{shape_name}': {', '.join(adjustments)}")
                
                except Exception as e:
                    print(f"    warning: could not validate shape on slide {slide_num}: {e}")
        
        if issues_found == 0:
            print(f"  all shapes within bounds")
        else:
            print(f"  fixed {issues_fixed} shapes that were out of bounds")
    
    def _enforce_text_fit_by_measurement(self, presentation):
        """
        FINAL SAFETY NET: Use actual text frame measurements to detect overflow.
        Every text box MUST fit within its template boundaries.
        
        NOTE: This is a last-resort enforcement. The AI service should have already:
        1. Generated content within capacity limits
        2. Created continuation slides for overflow content
        3. Used bullet points without markers
        
        If this enforcement truncates text, it means the AI's capacity estimation was off.
        The logs will show exactly what was truncated.
        """
        total_shortened = 0
        total_checked = 0
        overflow_reports = []
        
        print(f"  FINAL SAFETY NET: checking all text boxes for overflow...")
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_shortened = 0
            
            for shape in slide.shapes:
                try:
                    # Only process text frames
                    if not hasattr(shape, 'text_frame'):
                        continue
                    
                    text_frame = shape.text_frame
                    if not text_frame.text.strip():
                        continue
                    
                    total_checked += 1
                    
                    # Get original text
                    original_text = text_frame.text
                    original_len = len(original_text)
                    
                    # Get box metrics for debugging
                    box_info = self._get_box_info(shape, text_frame)
                    shape_name = getattr(shape, 'name', 'Unknown')
                    
                    # CRITICAL: Detect if text overflows template boundaries
                    needs_shortening = self._detect_text_overflow(shape, text_frame)
                    
                    if needs_shortening:
                        print(f"    ⚠️  slide {slide_num} '{shape_name}': OVERFLOW DETECTED")
                        print(f"        text: {original_len} chars, '{original_text[:60]}...'")
                        print(f"        box: {box_info}")
                        
                        # Iteratively shorten until it PHYSICALLY fits
                        fitted_text = self._iteratively_fit_text(shape, text_frame, original_text)
                        
                        if fitted_text != original_text:
                            total_shortened += 1
                            slide_shortened += 1
                            overflow_content = original_text[len(fitted_text):]
                            text_frame.text = fitted_text
                            reduction_pct = int((1 - len(fitted_text)/original_len) * 100)
                            
                            # Store overflow report
                            overflow_reports.append({
                                'slide_number': slide_num,
                                'shape_name': shape_name,
                                'box_info': box_info,
                                'original_length': original_len,
                                'fitted_length': len(fitted_text),
                                'overflow_content': overflow_content[:100] + '...' if len(overflow_content) > 100 else overflow_content
                            })
                            
                            print(f"    CORRECTED: {original_len} → {len(fitted_text)} chars (-{reduction_pct}%)")
                            print(f"    LOST CONTENT: '{overflow_content[:80]}...'")
                    else:
                        # Log that this box was checked and is OK
                        if original_len > 200:  # Only log longer text that passed
                            print(f"    slide {slide_num} '{shape_name}': {original_len} chars fits OK")
                
                except Exception as e:
                    print(f"    error processing text on slide {slide_num}: {e}")
                    import traceback
                    traceback.print_exc()
            
            if slide_shortened > 0:
                print(f"  📊 slide {slide_num}: corrected {slide_shortened} text boxes")
        
        print(f"\n  ENFORCEMENT COMPLETE:")
        print(f"     checked: {total_checked} text boxes")
        print(f"     corrected: {total_shortened} overflow violations")
        print(f"     status: {'all text within boundaries' if total_shortened == 0 else 'overflows fixed (content truncated)'}")
        
        if overflow_reports:
            print(f"\n  OVERFLOW REPORT ({len(overflow_reports)} truncations):")
            for report in overflow_reports:
                print(f"     • Slide {report['slide_number']} ({report['shape_name']}): {report['original_length']} → {report['fitted_length']} chars")
                print(f"       Lost: '{report['overflow_content']}'")
        
        return overflow_reports
    
    def _detect_text_overflow(self, shape, text_frame):
        """
        Detect text overflow by attempting to fit text and checking if font size would change.
        This uses python-pptx's actual text fitting logic for accurate detection.
        """
        from pptx.enum.text import MSO_AUTO_SIZE
        
        try:
            # Store ALL original font sizes from all runs
            original_font_sizes = []
            if text_frame.paragraphs:
                for para in text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            original_font_sizes.append(run.font.size)
            
            # If no explicit font sizes found, use default
            if not original_font_sizes:
                original_font_sizes = [914400 * 18 / 72]  # 18pt default
            
            # Use the largest font size as reference
            max_original_font_size = max(original_font_sizes)
            
            # Store original auto_size setting
            original_auto_size = text_frame.auto_size
            
            # Attempt to fit text to shape - this will reduce font size if overflow
            text_frame.word_wrap = True
            
            try:
                text_frame.fit_text(max_size=max_original_font_size)
            except AttributeError:
                # fit_text not available in this version, use fallback
                return self._detect_text_overflow_fallback(shape, text_frame)
            
            # Check if ANY font size was reduced (indicates overflow)
            fitted_font_sizes = []
            if text_frame.paragraphs:
                for para in text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            fitted_font_sizes.append(run.font.size)
            
            # Restore original auto_size
            text_frame.auto_size = original_auto_size
            
            # If any font size was reduced by more than 1%, text overflows
            overflow_detected = False
            if fitted_font_sizes:
                min_fitted = min(fitted_font_sizes)
                if min_fitted < max_original_font_size * 0.99:
                    shape_name = getattr(shape, 'name', 'Unknown')
                    reduction_pct = int((1 - min_fitted/max_original_font_size) * 100)
                    print(f"      overflow detected: font would need to shrink by {reduction_pct}% to fit")
                    overflow_detected = True
            
            # CRITICAL: Restore ALL original font sizes
            run_idx = 0
            for para in text_frame.paragraphs:
                for run in para.runs:
                    if run_idx < len(original_font_sizes):
                        run.font.size = original_font_sizes[run_idx]
                        run_idx += 1
            
            return overflow_detected
        
        except Exception as e:
            # If fit_text fails or isn't supported, fall back to heuristic
            print(f"      warning: fit_text detection failed ({e}), using fallback")
            import traceback
            traceback.print_exc()
            return self._detect_text_overflow_fallback(shape, text_frame)
    
    def _get_box_info(self, shape, text_frame):
        """Get human-readable box dimensions for debugging."""
        try:
            width_in = shape.width / 914400
            height_in = shape.height / 914400
            
            margin_left = getattr(text_frame, 'margin_left', 91440) / 914400
            margin_right = getattr(text_frame, 'margin_right', 91440) / 914400
            margin_top = getattr(text_frame, 'margin_top', 45720) / 914400
            margin_bottom = getattr(text_frame, 'margin_bottom', 45720) / 914400
            
            usable_width = width_in - margin_left - margin_right
            usable_height = height_in - margin_top - margin_bottom
            
            font_size = None
            if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                if text_frame.paragraphs[0].runs[0].font.size:
                    font_size = int(text_frame.paragraphs[0].runs[0].font.size / 914400 * 72)
            
            return f"{usable_width:.1f}\"×{usable_height:.1f}\" (font: {font_size}pt)"
        except:
            return "unknown dimensions"
    
    def _detect_text_overflow_fallback(self, shape, text_frame):
        """
        Fallback overflow detection using geometry calculations.
        Used if fit_text() method fails.
        """
        try:
            box_height = shape.height
            box_width = shape.width
            
            margin_left = getattr(text_frame, 'margin_left', 91440)
            margin_right = getattr(text_frame, 'margin_right', 91440)
            margin_top = getattr(text_frame, 'margin_top', 45720)
            margin_bottom = getattr(text_frame, 'margin_bottom', 45720)
            
            usable_height = box_height - margin_top - margin_bottom
            usable_width = box_width - margin_left - margin_right
            
            font_size_emu = 914400 * 18 / 72
            if text_frame.paragraphs:
                first_para = text_frame.paragraphs[0]
                if first_para.runs and first_para.runs[0].font.size:
                    font_size_emu = first_para.runs[0].font.size
            
            actual_line_count = self._count_wrapped_lines(text_frame.text, usable_width, font_size_emu)
            line_height = font_size_emu * 1.5
            height_needed = actual_line_count * line_height * 1.1
            
            return height_needed > usable_height
        
        except Exception:
            return len(text_frame.text) > 300
    
    def _count_wrapped_lines(self, text, usable_width_emu, font_size_emu):
        """
        Count actual number of lines after word wrapping.
        Uses realistic character width to simulate text rendering.
        """
        if not text.strip():
            return 0
        
        # Average character width in EMUs (proportional to font size)
        # Use 0.6 * font_size as a realistic average for proportional fonts
        char_width = font_size_emu * 0.6
        
        # Calculate chars that fit per line
        chars_per_line = max(1, int(usable_width_emu / char_width))
        
        # Split text into paragraphs (handle explicit newlines)
        paragraphs = text.split('\n')
        
        total_lines = 0
        for para in paragraphs:
            if not para.strip():
                total_lines += 1  # blank line
                continue
            
            # Wrap this paragraph into lines
            words = para.split()
            if not words:
                total_lines += 1
                continue
            
            current_line_length = 0
            line_count = 1
            
            for word in words:
                word_length = len(word) + 1  # +1 for space
                
                if current_line_length + word_length > chars_per_line:
                    # Word doesn't fit, wrap to next line
                    line_count += 1
                    current_line_length = word_length
                else:
                    current_line_length += word_length
            
            total_lines += line_count
        
        return total_lines
    
    def _iteratively_fit_text(self, shape, text_frame, original_text):
        """
        Iteratively shorten text until it PHYSICALLY fits within the box boundaries.
        Uses binary search with ACTUAL overflow detection (fit_text method).
        """
        if not original_text.strip():
            return original_text
        
        print(f"        fitting {len(original_text)} chars...")
        
        # Binary search to find maximum text that fits
        min_chars = 0
        max_chars = len(original_text)
        best_fit = ""
        
        # Start by testing if full text fits
        text_frame.text = original_text
        if not self._detect_text_overflow(shape, text_frame):
            print(f"        full text fits!")
            return original_text  # Perfect, no shortening needed
        
        # Binary search for maximum fitting length (max 20 iterations for precision)
        iteration_count = 0
        for iteration in range(20):
            iteration_count += 1
            
            if max_chars - min_chars <= 5:  # Very close
                break
            
            mid_chars = (min_chars + max_chars) // 2
            if mid_chars <= 0:
                break
            
            # Try text at this length
            test_text = self._truncate_at_word_boundary(original_text, mid_chars)
            text_frame.text = test_text
            
            # Check if it ACTUALLY overflows using fit_text
            overflows = self._detect_text_overflow(shape, text_frame)
            
            if overflows:
                # Still too much, try less
                max_chars = mid_chars
                print(f"        iter {iteration+1}: {len(test_text)} chars - still overflows")
            else:
                # Fits! Save this and try more
                best_fit = test_text
                min_chars = mid_chars
                print(f"        iter {iteration+1}: {len(test_text)} chars - fits!")
        
        print(f"        converged after {iteration_count} iterations")
        
        # If we found something that fits, use it
        if best_fit:
            # Add ellipsis if truncated
            if len(best_fit) < len(original_text):
                best_fit = best_fit.rstrip() + '...'
            print(f"        final: {len(best_fit)} chars")
            return best_fit
        
        # If nothing fits, use a minimal fallback
        print(f"        extreme overflow - using minimal fallback")
        fallback = self._truncate_at_word_boundary(original_text, 30)
        return fallback + '...'
    
    def _truncate_at_word_boundary(self, text, max_chars):
        """truncate text at word boundary near max_chars."""
        if len(text) <= max_chars:
            return text
        
        truncated = text[:max_chars]
        
        # find last space
        last_space = truncated.rfind(' ')
        if last_space > max_chars * 0.8:  # if we found a space reasonably close
            truncated = truncated[:last_space]
        
        return truncated.rstrip()
    
    def _validate_no_overflow(self, presentation):
        """
        Final validation pass: ensure NO text boxes overflow their boundaries.
        Returns count of violations found.
        """
        violations = 0
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            for shape in slide.shapes:
                try:
                    if not hasattr(shape, 'text_frame'):
                        continue
                    
                    text_frame = shape.text_frame
                    if not text_frame.text.strip():
                        continue
                    
                    # Check if this box overflows
                    if self._detect_text_overflow(shape, text_frame):
                        violations += 1
                        shape_name = getattr(shape, 'name', 'Unknown')
                        print(f"    ❌ VIOLATION: slide {slide_num}, shape '{shape_name}' still overflows!")
                
                except Exception as e:
                    print(f"    warning: validation error on slide {slide_num}: {e}")
        
        if violations == 0:
            print(f"  VALIDATION PASSED: Zero overflow violations detected")
        else:
            print(f"  VALIDATION FAILED: {violations} boxes still overflow boundaries")
        
        return violations
    
    def generate_deck(self, slide_specs, custom_theme=None):
        if not self.template_path or not os.path.exists(self.template_path):
            raise ValueError('No template presentation available')
        
        if self.stored_rules is None:
            raise ValueError('No rules stored')
        
        if not isinstance(slide_specs, list):
            raise ValueError(f'slide_specs must be a list, got {type(slide_specs)}')
        
        # Load BOTH the source (for cloning) and target (for output)
        source_prs = Presentation(self.template_path)
        target_prs = Presentation(self.template_path)
        
        # Clear all slides from target
        for i in range(len(target_prs.slides) - 1, -1, -1):
            rId = target_prs.slides._sldIdLst[i].rId
            target_prs.part.drop_rel(rId)
            del target_prs.slides._sldIdLst[i]
        
        print(f"\ncloning {len(slide_specs)} slides from template")
        
        # Apply custom theme overrides if provided
        if custom_theme:
            print(f"\napplying custom theme overrides...")
            self._apply_theme_overrides(custom_theme)
        
        layouts = self.stored_rules.get('layouts', [])
        layout_map = {layout['name']: layout for layout in layouts}
        
        # Import the cloning utility
        from utils.slide_cloner import clone_slide_with_content
        
        for i, spec in enumerate(slide_specs):
            if not isinstance(spec, dict):
                raise ValueError(f'slide_spec {i} must be a dict, got {type(spec)}: {spec}')
            
            layout_name = spec.get('layout_name')
            placeholders_data = spec.get('placeholders', [])
            
            print(f"\ngenerating slide {i + 1}: {layout_name}")
            print(f"content items: {len(placeholders_data)}")
            
            # Debug: show what content we're passing
            for ph in placeholders_data:
                ph_idx = ph.get('idx', '?')
                ph_type = ph.get('type', '?')
                if ph_type == 'text':
                    content_preview = ph.get('content', '')[:50]
                    print(f"  placeholder idx={ph_idx}, type={ph_type}, content='{content_preview}...'")
                else:
                    print(f"  placeholder idx={ph_idx}, type={ph_type}")
            
            if layout_name not in layout_map:
                print(f"  error: layout '{layout_name}' not found, skipping")
                continue
            
            layout_template = layout_map[layout_name]
            
            # Get the source slide index (0-based) from the layout
            source_slide_index = layout_template.get('layout_index', layout_template.get('slide_number', 1) - 1)
            
            try:
                # Clone the slide from source template, filling placeholders with new content
                clone_slide_with_content(
                    source_prs,
                    source_slide_index,
                    target_prs,
                    placeholders_data,
                    self.uploaded_images,
                    self.image_order
                )
                print(f"  ✓ slide cloned successfully from template slide {source_slide_index + 1}")
            except Exception as e:
                print(f"  error cloning slide: {e}")
                import traceback
                traceback.print_exc()
                continue
        
        # CRITICAL: Validate and fix all shape positions to ensure they're within slide bounds
        print(f"\nvalidating shape positions...")
        self._validate_and_fix_shape_positions(target_prs)
        
        # CRITICAL: Final safety net - enforce text fitting using actual measurements
        print(f"\nFINAL SAFETY NET: Overflow Enforcement...")
        overflow_reports = self._enforce_text_fit_by_measurement(target_prs)
        
        # FINAL VALIDATION: Double-check no overflows remain
        print(f"\nFINAL VALIDATION...")
        violations = self._validate_no_overflow(target_prs)
        if violations > 0:
            raise ValueError(f"CRITICAL: {violations} text boxes still overflow after enforcement! Generation failed.")
        
        # Warn if content was truncated (should rarely happen with good AI chunking)
        if overflow_reports:
            print(f"\nWARNING: {len(overflow_reports)} text boxes required truncation.")
            print(f"    This indicates AI capacity estimation was off. Content may be lost.")
        
        output = BytesIO()
        target_prs.save(output)
        output.seek(0)
        
        return base64.b64encode(output.getvalue()).decode('utf-8')
