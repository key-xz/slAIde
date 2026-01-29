import base64
from io import BytesIO
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor


def apply_text_frame_properties(text_frame, props):
    if not props:
        return
    
    applied_props = []
    try:
        if 'margin_left' in props and props['margin_left'] is not None:
            text_frame.margin_left = props['margin_left']
            applied_props.append(f"margin_left={props['margin_left']}")
        if 'margin_right' in props and props['margin_right'] is not None:
            text_frame.margin_right = props['margin_right']
            applied_props.append(f"margin_right={props['margin_right']}")
        if 'margin_top' in props and props['margin_top'] is not None:
            text_frame.margin_top = props['margin_top']
            applied_props.append(f"margin_top={props['margin_top']}")
        if 'margin_bottom' in props and props['margin_bottom'] is not None:
            text_frame.margin_bottom = props['margin_bottom']
            applied_props.append(f"margin_bottom={props['margin_bottom']}")
        if 'word_wrap' in props and props['word_wrap'] is not None:
            text_frame.word_wrap = props['word_wrap']
            applied_props.append(f"word_wrap={props['word_wrap']}")
        if 'vertical_anchor' in props and props['vertical_anchor']:
            anchor_str = props['vertical_anchor']
            if 'TOP' in anchor_str:
                text_frame.vertical_anchor = MSO_ANCHOR.TOP
                applied_props.append("anchor=TOP")
            elif 'MIDDLE' in anchor_str:
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                applied_props.append("anchor=MIDDLE")
            elif 'BOTTOM' in anchor_str:
                text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
                applied_props.append("anchor=BOTTOM")
        
        if applied_props:
            print(f"      ✓ applied text frame: {', '.join(applied_props)}")
    except Exception as e:
        print(f"    warning: failed to apply text frame props: {e}")
        import traceback
        traceback.print_exc()


def apply_paragraph_properties(paragraph, props):
    if not props:
        return
    
    applied_props = []
    try:
        if 'alignment' in props and props['alignment']:
            align_str = props['alignment']
            if 'LEFT' in align_str:
                paragraph.alignment = PP_ALIGN.LEFT
                applied_props.append("align=LEFT")
            elif 'CENTER' in align_str:
                paragraph.alignment = PP_ALIGN.CENTER
                applied_props.append("align=CENTER")
            elif 'RIGHT' in align_str:
                paragraph.alignment = PP_ALIGN.RIGHT
                applied_props.append("align=RIGHT")
            elif 'JUSTIFY' in align_str:
                paragraph.alignment = PP_ALIGN.JUSTIFY
                applied_props.append("align=JUSTIFY")
        
        if 'line_spacing' in props and props['line_spacing']:
            paragraph.line_spacing = props['line_spacing']
            applied_props.append(f"line_spacing={props['line_spacing']}")
        
        if 'space_before' in props and props['space_before']:
            paragraph.space_before = Pt(props['space_before'])
            applied_props.append(f"space_before={props['space_before']}pt")
        
        if 'space_after' in props and props['space_after']:
            paragraph.space_after = Pt(props['space_after'])
            applied_props.append(f"space_after={props['space_after']}pt")
        
        if 'level' in props:
            paragraph.level = props['level']
            applied_props.append(f"level={props['level']}")
        
        if applied_props:
            print(f"      ✓ applied paragraph: {', '.join(applied_props)}")
    except Exception as e:
        print(f"    warning: failed to apply paragraph props: {e}")
        import traceback
        traceback.print_exc()


def apply_font_properties(font, props):
    if not props:
        print(f"      ✗ no font props provided")
        return
    
    print(f"      → attempting to apply font props: {props}")
    
    try:
        if 'name' in props:
            if props['name']:
                font.name = props['name']
                print(f"      ✓ applied font name: {props['name']}")
            else:
                print(f"      ✗ font name is None/empty")
        else:
            print(f"      ✗ no 'name' key in font props")
        
        if 'size' in props:
            if props['size']:
                font.size = Pt(props['size'])
                print(f"      ✓ applied font size: {props['size']}pt")
            else:
                print(f"      ✗ font size is None/empty")
        else:
            print(f"      ✗ no 'size' key in font props")
        
        if 'bold' in props and props['bold'] is not None:
            font.bold = props['bold']
            print(f"      ✓ applied font bold: {props['bold']}")
        
        if 'italic' in props and props['italic'] is not None:
            font.italic = props['italic']
            print(f"      ✓ applied font italic: {props['italic']}")
        
        if 'underline' in props and props['underline'] is not None:
            font.underline = props['underline']
            print(f"      ✓ applied font underline: {props['underline']}")
        
        if 'color' in props:
            color_info = props['color']
            
            # Check if we have a valid RGB color
            if color_info and 'rgb' in color_info and color_info['rgb']:
                rgb_str = color_info['rgb']
                if rgb_str and len(rgb_str) >= 6:
                    rgb_hex = rgb_str[:6]
                    r = int(rgb_hex[0:2], 16)
                    g = int(rgb_hex[2:4], 16)
                    b = int(rgb_hex[4:6], 16)
                    font.color.rgb = RGBColor(r, g, b)
                    print(f"      ✓ applied font color: #{rgb_hex}")
            # If color type is 'None' or invalid, apply default black
            elif color_info and color_info.get('type') == 'None':
                font.color.rgb = RGBColor(0, 0, 0)  # Black
                print(f"      ✓ applied default black color (original was None)")
            elif 'theme_color' in color_info and color_info['theme_color']:
                print(f"      → theme color detected: {color_info['theme_color']} (preserved)")
    except Exception as e:
        print(f"    warning: failed to apply font props: {e}")
        import traceback
        traceback.print_exc()


def apply_fill_properties(shape, props):
    if not props:
        return
    
    try:
        if 'type' in props and 'SOLID' in props['type']:
            shape.fill.solid()
            
            if 'fore_color' in props:
                fore_color = props['fore_color']
                if 'rgb' in fore_color:
                    rgb_str = fore_color['rgb']
                    if rgb_str and len(rgb_str) >= 6:
                        r = int(rgb_str[0:2], 16)
                        g = int(rgb_str[2:4], 16)
                        b = int(rgb_str[4:6], 16)
                        shape.fill.fore_color.rgb = RGBColor(r, g, b)
    except Exception as e:
        print(f"    warning: failed to apply fill props: {e}")


def apply_line_properties(shape, props):
    if not props:
        return
    
    try:
        if 'width' in props:
            shape.line.width = props['width']
        
        if 'color' in props:
            line_color = props['color']
            if 'rgb' in line_color:
                rgb_str = line_color['rgb']
                if rgb_str and len(rgb_str) >= 6:
                    r = int(rgb_str[0:2], 16)
                    g = int(rgb_str[2:4], 16)
                    b = int(rgb_str[4:6], 16)
                    shape.line.color.rgb = RGBColor(r, g, b)
    except Exception as e:
        print(f"    warning: failed to apply line props: {e}")


def create_text_shape(slide, shape_props, content):
    pos = shape_props['position']
    left = pos['left']
    top = pos['top']
    width = pos['width']
    height = pos['height']
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    print(f"      → creating text box at ({left}, {top}) size ({width}x{height})")
    print(f"      → content to insert: '{content[:100]}...' ({len(content)} chars)")
    
    if 'text_frame_props' in shape_props:
        apply_text_frame_properties(text_frame, shape_props['text_frame_props'])
    
    text_frame.clear()
    p = text_frame.paragraphs[0]
    
    if 'paragraph_props' in shape_props:
        apply_paragraph_properties(p, shape_props['paragraph_props'])
    
    run = p.add_run()
    run.text = content
    
    print(f"      → text set on run: '{run.text[:50]}...'")
    print(f"      → text frame has text: '{text_frame.text[:50] if text_frame.text else 'EMPTY'}...'")
    
    if 'font_props' in shape_props:
        print(f"      → applying font properties: {shape_props['font_props']}")
        apply_font_properties(run.font, shape_props['font_props'])
    else:
        print(f"      → no font properties found in template")
    
    if 'fill_props' in shape_props:
        apply_fill_properties(textbox, shape_props['fill_props'])
    
    if 'line_props' in shape_props:
        apply_line_properties(textbox, shape_props['line_props'])
    
    if 'rotation' in shape_props:
        textbox.rotation = shape_props['rotation']
    
    print(f"      ✅ textbox complete: has_text={bool(text_frame.text)}, text_length={len(text_frame.text or '')}")
    return textbox


def create_image_shape(slide, shape_props, image_data):
    pos = shape_props['position']
    left = pos['left']
    top = pos['top']
    width = pos['width']
    height = pos['height']
    
    if ',' in image_data:
        image_data = image_data.split(',')[1]
    
    image_bytes = base64.b64decode(image_data)
    image_stream = BytesIO(image_bytes)
    
    picture = slide.shapes.add_picture(image_stream, left, top, width, height)
    
    if 'rotation' in shape_props:
        picture.rotation = shape_props['rotation']
    
    return picture


def create_static_shape(slide, shape_props):
    shape_type_str = shape_props.get('shape_type', '')
    shape_name = shape_props.get('name', 'unnamed')
    pos = shape_props['position']
    left = pos['left']
    top = pos['top']
    width = pos['width']
    height = pos['height']
    
    # Handle design images (logos, backgrounds, decorative elements)
    if 'PICTURE' in shape_type_str:
        if shape_props.get('is_design_image') and 'image_data' in shape_props and shape_props['image_data']:
            try:
                image_data = shape_props['image_data']
                image_bytes = base64.b64decode(image_data)
                image_stream = BytesIO(image_bytes)
                
                picture = slide.shapes.add_picture(image_stream, left, top, width, height)
                
                if 'rotation' in shape_props:
                    picture.rotation = shape_props['rotation']
                
                print(f"      ✓ preserved design image '{shape_name}'")
                return picture
            except Exception as e:
                print(f"      ✗ failed to add design image '{shape_name}': {e}")
                return None
        else:
            print(f"      ⊘ skipping picture '{shape_name}' (no data or not design image)")
            return None
    
    # Handle text boxes (static text content)
    if 'TEXT_BOX' in shape_type_str:
        try:
            textbox = slide.shapes.add_textbox(left, top, width, height)
            
            if 'text_frame_props' in shape_props:
                apply_text_frame_properties(textbox.text_frame, shape_props['text_frame_props'])
            
            # Try to preserve any static text content
            # (placeholder text won't be here since placeholders are handled separately)
            
            if 'fill_props' in shape_props:
                apply_fill_properties(textbox, shape_props['fill_props'])
            if 'line_props' in shape_props:
                apply_line_properties(textbox, shape_props['line_props'])
            if 'rotation' in shape_props:
                textbox.rotation = shape_props['rotation']
            
            print(f"      ✓ preserved text box '{shape_name}'")
            return textbox
        except Exception as e:
            print(f"      ✗ failed to add text box '{shape_name}': {e}")
            return None
    
    # Handle auto shapes (rectangles, circles, etc.)
    if 'AUTO_SHAPE' in shape_type_str or 'SHAPE' in shape_type_str:
        try:
            # Default to rectangle for now - would need shape type mapping for specific shapes
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
            
            if 'fill_props' in shape_props:
                apply_fill_properties(shape, shape_props['fill_props'])
            if 'line_props' in shape_props:
                apply_line_properties(shape, shape_props['line_props'])
            if 'rotation' in shape_props:
                shape.rotation = shape_props['rotation']
            
            print(f"      ✓ preserved shape '{shape_name}'")
            return shape
        except Exception as e:
            print(f"      ✗ failed to add shape '{shape_name}': {e}")
            return None
    
    # For unsupported types, log and skip
    print(f"      ⊘ skipping unsupported shape type '{shape_name}' [{shape_type_str}]")
    return None


def generate_slide_from_template(presentation, layout_template, placeholder_contents, uploaded_images):
    """
    Generate a slide by recreating it from our extracted template.
    
    CRITICAL: We use a BLANK layout and rebuild everything ourselves because:
    - We extracted from actual slides (presentation.slides[N])
    - Not from master layouts (presentation.slide_layouts[N])
    - Using master layouts gives us wrong/missing design elements
    """
    # Use the blankest layout available (usually index 6 is "Blank" but not always)
    # Try to find a blank layout, otherwise use the last one
    slide_layout = None
    for layout in presentation.slide_layouts:
        if 'blank' in layout.name.lower():
            slide_layout = layout
            break
    
    if slide_layout is None:
        # Fallback: use the layout with the fewest placeholders (most blank)
        slide_layout = min(presentation.slide_layouts, 
                          key=lambda l: len([s for s in l.placeholders]))
    
    print(f"    using layout: '{slide_layout.name}' to build from scratch")
    slide = presentation.slides.add_slide(slide_layout)
    
    # CRITICAL: Remove ALL existing shapes from the blank layout
    # We'll rebuild everything from our extracted template
    shapes_to_remove = list(slide.shapes)
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
    
    print(f"    removed {len(shapes_to_remove)} default shapes, starting fresh")
    
    # Since we removed all shapes, we need to create everything from scratch
    print(f"    recreating {len(layout_template['placeholders'])} placeholders from template...")
    print(f"    received {len(placeholder_contents)} content items")
    
    content_map = {item['idx']: item for item in placeholder_contents}
    
    # Debug: show what content we have
    print(f"    content_map keys: {list(content_map.keys())}")
    for idx, item in content_map.items():
        item_type = item.get('type', 'unknown')
        if item_type == 'text':
            preview = item.get('content', '')[:30]
            print(f"      idx={idx}: type={item_type}, preview='{preview}...'")
        else:
            print(f"      idx={idx}: type={item_type}")
    
    # Create all placeholders from our extracted template
    for placeholder_def in layout_template['placeholders']:
        idx = placeholder_def['idx']
        ph_type = placeholder_def['type']
        shape_props = placeholder_def['properties']
        
        print(f"    processing placeholder idx={idx}, type={ph_type}")
        
        if idx not in content_map:
            print(f"    ⚠️  WARNING: no content provided for placeholder idx={idx}")
            continue
        
        content_item = content_map[idx]
        content_type = content_item.get('type')
        
        print(f"      content_item type: {content_type}")
        
        if ph_type == 'text' and content_type == 'text':
            content = content_item.get('content', '')
            print(f"      ✓ creating TEXT shape with {len(content)} chars: '{content[:50]}...'")
            result = create_text_shape(slide, shape_props, content)
            if result:
                print(f"      ✅ text shape created successfully")
            else:
                print(f"      ❌ text shape creation returned None")
        elif ph_type != content_type:
            print(f"      ❌ TYPE MISMATCH: placeholder expects {ph_type} but content is {content_type}")
        
        elif ph_type == 'image' and content_item.get('type') == 'image':
            image_index = content_item.get('image_index')
            if image_index is not None:
                image_filenames = list(uploaded_images.keys())
                if 0 <= image_index < len(image_filenames):
                    filename = image_filenames[image_index]
                    image_data = uploaded_images[filename]
                    print(f"    creating image placeholder idx={idx}: {filename}")
                    create_image_shape(slide, shape_props, image_data)
                else:
                    print(f"    error: image_index {image_index} out of range")
            else:
                print(f"    error: no image_index for image placeholder")
    
    # Recreate all static shapes (design elements, backgrounds, etc.)
    static_shapes = layout_template.get('shapes', [])
    if static_shapes:
        print(f"    recreating {len(static_shapes)} static shapes...")
        for static_shape_props in static_shapes:
            create_static_shape(slide, static_shape_props)
    else:
        print(f"    no static shapes to recreate")
    
    # Summary
    final_shape_count = len(list(slide.shapes))
    print(f"    ✅ slide complete: {final_shape_count} total shapes")
    print(f"       ({len(layout_template['placeholders'])} placeholders + {len(static_shapes)} static)")
    
    return slide
