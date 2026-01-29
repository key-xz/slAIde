from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def extract_placeholder_complete_styling(placeholder):
    styling = {
        'position': {
            'left': placeholder.left,
            'top': placeholder.top,
            'width': placeholder.width,
            'height': placeholder.height
        }
    }
    
    if hasattr(placeholder, 'text_frame'):
        text_frame = placeholder.text_frame
        
        text_frame_props = {}
        try:
            if hasattr(text_frame, 'margin_left'):
                text_frame_props['margin_left'] = text_frame.margin_left
            if hasattr(text_frame, 'margin_right'):
                text_frame_props['margin_right'] = text_frame.margin_right
            if hasattr(text_frame, 'margin_top'):
                text_frame_props['margin_top'] = text_frame.margin_top
            if hasattr(text_frame, 'margin_bottom'):
                text_frame_props['margin_bottom'] = text_frame.margin_bottom
            if hasattr(text_frame, 'word_wrap'):
                text_frame_props['word_wrap'] = text_frame.word_wrap
            if hasattr(text_frame, 'vertical_anchor'):
                text_frame_props['vertical_anchor'] = str(text_frame.vertical_anchor)
        except:
            pass
        
        if text_frame_props:
            styling['text_frame'] = text_frame_props
        
        if text_frame.paragraphs:
            para = text_frame.paragraphs[0]
            
            para_props = {}
            try:
                if hasattr(para, 'alignment') and para.alignment:
                    para_props['alignment'] = str(para.alignment)
                if hasattr(para, 'line_spacing'):
                    para_props['line_spacing'] = float(para.line_spacing) if para.line_spacing else None
                if hasattr(para, 'space_before'):
                    para_props['space_before'] = para.space_before.pt if para.space_before else None
                if hasattr(para, 'space_after'):
                    para_props['space_after'] = para.space_after.pt if para.space_after else None
                if hasattr(para, 'level'):
                    para_props['level'] = para.level
            except:
                pass
            
            if para_props:
                styling['paragraph'] = para_props
            
            if para.runs:
                run = para.runs[0]
                font = run.font
                
                font_props = {}
                try:
                    if font.name:
                        font_props['name'] = font.name
                    if font.size:
                        font_props['size'] = font.size.pt
                    if font.bold is not None:
                        font_props['bold'] = font.bold
                    if font.italic is not None:
                        font_props['italic'] = font.italic
                    if font.underline is not None:
                        font_props['underline'] = bool(font.underline)
                    
                    if hasattr(font, 'color') and font.color:
                        color_info = {}
                        try:
                            if hasattr(font.color, 'type'):
                                color_info['type'] = str(font.color.type)
                            if hasattr(font.color, 'rgb') and font.color.rgb:
                                rgb_obj = font.color.rgb
                                rgb_hex = format(rgb_obj[0], '02X') + format(rgb_obj[1], '02X') + format(rgb_obj[2], '02X')
                                color_info['rgb'] = rgb_hex
                            if hasattr(font.color, 'theme_color') and font.color.theme_color:
                                color_info['theme_color'] = str(font.color.theme_color)
                            if hasattr(font.color, 'brightness') and font.color.brightness is not None:
                                color_info['brightness'] = font.color.brightness
                        except:
                            pass
                        if color_info:
                            font_props['color'] = color_info
                except:
                    pass
                
                if font_props:
                    styling['font'] = font_props
    
    if hasattr(placeholder, 'fill'):
        try:
            fill_props = {}
            fill = placeholder.fill
            
            if hasattr(fill, 'type'):
                fill_props['type'] = str(fill.type)
            
            if hasattr(fill, 'fore_color') and fill.fore_color:
                fore_color = {}
                try:
                    if hasattr(fill.fore_color, 'rgb') and fill.fore_color.rgb:
                        fore_color['rgb'] = str(fill.fore_color.rgb)
                    if hasattr(fill.fore_color, 'theme_color') and fill.fore_color.theme_color:
                        fore_color['theme_color'] = str(fill.fore_color.theme_color)
                except:
                    pass
                if fore_color:
                    fill_props['fore_color'] = fore_color
            
            if fill_props:
                styling['fill'] = fill_props
        except:
            pass
    
    if hasattr(placeholder, 'line'):
        try:
            line_props = {}
            line = placeholder.line
            
            if hasattr(line, 'width') and line.width:
                line_props['width'] = line.width
            
            if hasattr(line, 'color') and line.color:
                line_color = {}
                try:
                    if hasattr(line.color, 'rgb') and line.color.rgb:
                        line_color['rgb'] = str(line.color.rgb)
                except:
                    pass
                if line_color:
                    line_props['color'] = line_color
            
            if line_props:
                styling['line'] = line_props
        except:
            pass
    
    return styling


def apply_placeholder_styling(placeholder, styling, new_text=None):
    if new_text and hasattr(placeholder, 'text_frame'):
        text_frame = placeholder.text_frame
        text_frame.clear()
        
        if 'text_frame' in styling:
            tf_props = styling['text_frame']
            try:
                if 'margin_left' in tf_props:
                    text_frame.margin_left = tf_props['margin_left']
                if 'margin_right' in tf_props:
                    text_frame.margin_right = tf_props['margin_right']
                if 'margin_top' in tf_props:
                    text_frame.margin_top = tf_props['margin_top']
                if 'margin_bottom' in tf_props:
                    text_frame.margin_bottom = tf_props['margin_bottom']
                if 'word_wrap' in tf_props:
                    text_frame.word_wrap = tf_props['word_wrap']
                if 'vertical_anchor' in tf_props:
                    anchor_str = tf_props['vertical_anchor']
                    if 'TOP' in anchor_str:
                        text_frame.vertical_anchor = MSO_ANCHOR.TOP
                    elif 'MIDDLE' in anchor_str:
                        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    elif 'BOTTOM' in anchor_str:
                        text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
            except:
                pass
        
        p = text_frame.paragraphs[0]
        
        if 'paragraph' in styling:
            para_props = styling['paragraph']
            try:
                if 'alignment' in para_props:
                    align_str = para_props['alignment']
                    if 'LEFT' in align_str:
                        p.alignment = PP_ALIGN.LEFT
                    elif 'CENTER' in align_str:
                        p.alignment = PP_ALIGN.CENTER
                    elif 'RIGHT' in align_str:
                        p.alignment = PP_ALIGN.RIGHT
                    elif 'JUSTIFY' in align_str:
                        p.alignment = PP_ALIGN.JUSTIFY
                
                if 'line_spacing' in para_props and para_props['line_spacing']:
                    p.line_spacing = para_props['line_spacing']
                
                if 'space_before' in para_props and para_props['space_before']:
                    p.space_before = Pt(para_props['space_before'])
                
                if 'space_after' in para_props and para_props['space_after']:
                    p.space_after = Pt(para_props['space_after'])
                
                if 'level' in para_props:
                    p.level = para_props['level']
            except:
                pass
        
        run = p.add_run()
        run.text = new_text
        
        if 'font' in styling:
            font_props = styling['font']
            try:
                if 'name' in font_props:
                    run.font.name = font_props['name']
                
                if 'size' in font_props:
                    run.font.size = Pt(font_props['size'])
                
                if 'bold' in font_props:
                    run.font.bold = font_props['bold']
                
                if 'italic' in font_props:
                    run.font.italic = font_props['italic']
                
                if 'underline' in font_props:
                    run.font.underline = font_props['underline']
                
                if 'color' in font_props:
                    color_info = font_props['color']
                    try:
                        if 'rgb' in color_info:
                            rgb_str = color_info['rgb']
                            if rgb_str and len(rgb_str) == 6:
                                from pptx.dml.color import RGBColor
                                r = int(rgb_str[0:2], 16)
                                g = int(rgb_str[2:4], 16)
                                b = int(rgb_str[4:6], 16)
                                run.font.color.rgb = RGBColor(r, g, b)
                    except:
                        pass
            except:
                pass
    
    if 'fill' in styling:
        try:
            fill_props = styling['fill']
            fill = placeholder.fill
            
            if 'type' in fill_props and 'SOLID' in fill_props['type']:
                fill.solid()
                
                if 'fore_color' in fill_props:
                    fore_color = fill_props['fore_color']
                    if 'rgb' in fore_color:
                        from pptx.dml.color import RGBColor
                        rgb_str = fore_color['rgb']
                        if rgb_str and len(rgb_str) == 6:
                            r = int(rgb_str[0:2], 16)
                            g = int(rgb_str[2:4], 16)
                            b = int(rgb_str[4:6], 16)
                            fill.fore_color.rgb = RGBColor(r, g, b)
        except:
            pass
    
    if 'line' in styling:
        try:
            line_props = styling['line']
            line = placeholder.line
            
            if 'width' in line_props:
                line.width = line_props['width']
            
            if 'color' in line_props:
                line_color = line_props['color']
                if 'rgb' in line_color:
                    from pptx.dml.color import RGBColor
                    rgb_str = line_color['rgb']
                    if rgb_str and len(rgb_str) == 6:
                        r = int(rgb_str[0:2], 16)
                        g = int(rgb_str[2:4], 16)
                        b = int(rgb_str[4:6], 16)
                        line.color.rgb = RGBColor(r, g, b)
        except:
            pass
