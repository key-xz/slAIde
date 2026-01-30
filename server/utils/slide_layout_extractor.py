import base64
import zipfile
from xml.etree import ElementTree as ET
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


def extract_shape_complete_properties(shape, default_font=None):
    """
    extracts all relevant properties from a shape.
    treats all text boxes and images as content areas for google slides compatibility.
    """
    if default_font is None:
        default_font = {'name': 'Calibri', 'size': 18}
    
    shape_type = shape.shape_type
    
    # check if it's a placeholder and its type
    is_placeholder = False
    real_idx = -1
    placeholder_type = None
    
    try:
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            is_placeholder = True
            if hasattr(shape, 'placeholder_format'):
                real_idx = shape.placeholder_format.idx
                # get the placeholder type to determine if it's for images
                if hasattr(shape.placeholder_format, 'type'):
                    placeholder_type = shape.placeholder_format.type
    except:
        pass
    
    base_props = {
        'placeholder_idx': real_idx,
        'shape_type': str(shape_type),
        'name': shape.name,
        'position': {
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height
        },
        'rotation': shape.rotation if hasattr(shape, 'rotation') else 0
    }
    
    # determine content type
    # check if it's an image placeholder (PP_PLACEHOLDER.PICTURE or similar)
    if placeholder_type is not None:
        # Check for image placeholder types (be flexible for different python-pptx versions)
        is_image_placeholder = False
        
        try:
            if placeholder_type == PP_PLACEHOLDER.PICTURE:
                is_image_placeholder = True
        except (AttributeError, TypeError):
            pass
        
        # CLIP_ART and MEDIA may not exist in all versions, check safely
        try:
            if hasattr(PP_PLACEHOLDER, 'CLIP_ART') and placeholder_type == PP_PLACEHOLDER.CLIP_ART:
                is_image_placeholder = True
        except (AttributeError, TypeError):
            pass
        
        try:
            if hasattr(PP_PLACEHOLDER, 'MEDIA') and placeholder_type == PP_PLACEHOLDER.MEDIA:
                is_image_placeholder = True
        except (AttributeError, TypeError):
            pass
        
        # Also check by string matching (more robust)
        type_str = str(placeholder_type).upper()
        if 'PICTURE' in type_str or 'IMAGE' in type_str or 'MEDIA' in type_str:
            is_image_placeholder = True
        
        if is_image_placeholder:
            base_props['content_type'] = 'image'
            base_props['is_placeholder'] = True
            print(f"      ✓ detected IMAGE placeholder: {shape.name} (type: {placeholder_type})")
            return base_props
    
    # check if it's an actual picture shape (existing design images)
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        base_props['content_type'] = 'image'
        
        # CRITICAL: Only mark as placeholder if it's actually a placeholder shape
        # If is_placeholder was True, we would have caught it above
        # So this is a design image (logo, background, etc.) that should be PRESERVED
        if is_placeholder:
            # This is an empty image placeholder (shouldn't normally reach here)
            base_props['is_placeholder'] = True
            print(f"      ✓ detected empty IMAGE placeholder: {shape.name}")
        else:
            # This is a design image that should be preserved as static
            base_props['is_placeholder'] = False
            base_props['is_design_image'] = True
            print(f"      ✓ detected DESIGN IMAGE (will be preserved): {shape.name}")
        
        try:
            base_props['image_data'] = base64.b64encode(shape.image.blob).decode('utf-8')
        except:
            pass
        return base_props
    
    if hasattr(shape, 'text_frame'):
        text_frame = shape.text_frame
        
        base_props['content_type'] = 'text'
        base_props['is_placeholder'] = True
        
        base_props['text_frame_props'] = {
            'margin_left': text_frame.margin_left if hasattr(text_frame, 'margin_left') else None,
            'margin_right': text_frame.margin_right if hasattr(text_frame, 'margin_right') else None,
            'margin_top': text_frame.margin_top if hasattr(text_frame, 'margin_top') else None,
            'margin_bottom': text_frame.margin_bottom if hasattr(text_frame, 'margin_bottom') else None,
            'word_wrap': text_frame.word_wrap if hasattr(text_frame, 'word_wrap') else None,
            'vertical_anchor': str(text_frame.vertical_anchor) if hasattr(text_frame, 'vertical_anchor') else None,
            'auto_size': str(text_frame.auto_size) if hasattr(text_frame, 'auto_size') else None
        }
        
        if text_frame.paragraphs:
            para = text_frame.paragraphs[0]
            print(f"      shape '{shape.name}' has {len(text_frame.paragraphs)} paragraph(s)")
            
            base_props['paragraph_props'] = {
                'alignment': str(para.alignment) if hasattr(para, 'alignment') and para.alignment else None,
                'line_spacing': float(para.line_spacing) if hasattr(para, 'line_spacing') and para.line_spacing else None,
                'space_before': para.space_before.pt if hasattr(para, 'space_before') and para.space_before else None,
                'space_after': para.space_after.pt if hasattr(para, 'space_after') and para.space_after else None,
                'level': para.level if hasattr(para, 'level') else 0
            }
            
            had_runs = len(para.runs) > 0
            if para.runs:
                print(f"      paragraph has {len(para.runs)} run(s)")
            else:
                print(f"      ✗ paragraph has NO runs - will try to add dummy text")
                try:
                    dummy_run = para.add_run()
                    dummy_run.text = "X"
                    print(f"      ✓ added dummy run to extract font properties")
                except Exception as e:
                    print(f"      ✗ failed to add dummy run: {e}")
            
            if para.runs:
                run = para.runs[0]
                font = run.font
                
                font_props = {}
                
                try:
                    name = font.name
                    print(f"        font.name value: {repr(name)} (type: {type(name).__name__})")
                    if name:
                        font_props['name'] = name
                        print(f"        ✓ extracted font.name: {name}")
                    else:
                        font_props['name'] = default_font['name']
                        print(f"        → using default font.name from master: {default_font['name']}")
                except Exception as e:
                    print(f"        ✗ error getting font.name: {e}")
                    font_props['name'] = default_font['name']
                
                try:
                    size = font.size
                    print(f"        font.size value: {repr(size)} (type: {type(size).__name__ if size else 'NoneType'})")
                    if size:
                        size_pt = size.pt
                        font_props['size'] = size_pt
                        print(f"        ✓ extracted font.size: {size_pt}pt")
                    else:
                        font_props['size'] = default_font['size']
                        print(f"        → using default font.size from master: {default_font['size']}pt")
                except Exception as e:
                    print(f"        ✗ error getting font.size: {e}")
                    font_props['size'] = default_font['size']
                
                if font.bold is not None:
                    font_props['bold'] = font.bold
                if font.italic is not None:
                    font_props['italic'] = font.italic
                if font.underline is not None:
                    font_props['underline'] = bool(font.underline)
                
                if hasattr(font, 'color') and font.color:
                    try:
                        color_info = {}
                        color_type_str = None
                        if hasattr(font.color, 'type'):
                            color_type_str = str(font.color.type)
                            color_info['type'] = color_type_str
                        
                        if hasattr(font.color, 'rgb') and font.color.rgb:
                            rgb_obj = font.color.rgb
                            rgb_hex = format(rgb_obj[0], '02X') + format(rgb_obj[1], '02X') + format(rgb_obj[2], '02X')
                            color_info['rgb'] = rgb_hex
                        if hasattr(font.color, 'theme_color') and font.color.theme_color:
                            color_info['theme_color'] = str(font.color.theme_color)
                        if hasattr(font.color, 'brightness') and font.color.brightness is not None:
                            color_info['brightness'] = font.color.brightness
                        
                        # Only store color if we have actual color data (RGB or theme)
                        # Don't store if it's just {'type': 'None'}
                        has_valid_color = 'rgb' in color_info or 'theme_color' in color_info
                        if not has_valid_color and color_type_str in ['None', 'MSO_COLOR_TYPE.RGB', '1']:
                            # No valid color, use default black
                            color_info['rgb'] = '000000'  # Black
                            print(f"      ✓ no color found, defaulting to black")
                        
                        if color_info and ('rgb' in color_info or 'theme_color' in color_info):
                            font_props['color'] = color_info
                    except Exception as e:
                        print(f"    warning: failed to extract font color: {e}")
                        pass
                
                if font_props:
                    base_props['font_props'] = font_props
                    font_summary = []
                    if 'name' in font_props:
                        font_summary.append(f"font={font_props['name']}")
                    if 'size' in font_props:
                        font_summary.append(f"size={font_props['size']}pt")
                    if 'color' in font_props and 'rgb' in font_props['color']:
                        font_summary.append(f"color=#{font_props['color']['rgb']}")
                    if font_summary:
                        print(f"      ✓ extracted font: {', '.join(font_summary)}")
                
                if not had_runs and para.runs:
                    try:
                        para.clear()
                        print(f"      ✓ cleaned up dummy run")
                    except:
                        pass
    else:
        base_props['content_type'] = 'static'
        base_props['is_placeholder'] = False
    
    if hasattr(shape, 'fill'):
        try:
            fill_props = {}
            fill = shape.fill
            
            if hasattr(fill, 'type'):
                fill_props['type'] = str(fill.type)
            
            if hasattr(fill, 'fore_color') and fill.fore_color:
                fore_color = {}
                try:
                    if hasattr(fill.fore_color, 'rgb') and fill.fore_color.rgb:
                        fore_color['rgb'] = str(fill.fore_color.rgb)
                    if hasattr(fill.fore_color, 'theme_color') and fill.fore_color.theme_color:
                        fore_color['theme_color'] = str(fill.fore_color.theme_color)
                    if hasattr(fill.fore_color, 'brightness'):
                        fore_color['brightness'] = fill.fore_color.brightness
                except:
                    pass
                if fore_color:
                    fill_props['fore_color'] = fore_color
            
            if fill_props:
                base_props['fill_props'] = fill_props
        except:
            pass
    
    if hasattr(shape, 'line'):
        try:
            line_props = {}
            line = shape.line
            
            if hasattr(line, 'width') and line.width:
                line_props['width'] = line.width
            
            if hasattr(line, 'color') and line.color:
                line_color = {}
                try:
                    if hasattr(line.color, 'rgb') and line.color.rgb:
                        line_color['rgb'] = str(line.color.rgb)
                    if hasattr(line.color, 'theme_color') and line.color.theme_color:
                        line_color['theme_color'] = str(line.color.theme_color)
                except:
                    pass
                if line_color:
                    line_props['color'] = line_color
            
            if line_props:
                base_props['line_props'] = line_props
        except:
            pass
    
    if hasattr(shape, 'shadow'):
        try:
            shadow_props = {}
            shadow = shape.shadow
            
            if hasattr(shadow, 'inherit') and shadow.inherit is not None:
                shadow_props['inherit'] = shadow.inherit
            
            if shadow_props:
                base_props['shadow_props'] = shadow_props
        except:
            pass
    
    return base_props


def extract_complete_theme_from_zip(pptx_path):
    """
    Extract ALL theme data from .pptx including fonts, colors, backgrounds, effects.
    Returns comprehensive theme information for preservation.
    """
    result = {
        'fonts': {
            'title': {'name': 'Calibri', 'size': 44},
            'body': {'name': 'Calibri', 'size': 18}
        },
        'color_scheme': {},
        'format_scheme': {},
        'backgrounds': [],
        'theme_raw': {}
    }
    
    try:
        print(f"      🎨 extracting comprehensive theme data from .pptx...")
        
        with zipfile.ZipFile(pptx_path, 'r') as zip_file:
            file_list = zip_file.namelist()
            print(f"      found {len(file_list)} files in .pptx")
            
            # Extract theme files
            theme_files = [f for f in file_list if 'theme' in f.lower() and f.endswith('.xml')]
            print(f"      found theme files: {theme_files}")
            
            # Extract slide master files (contain backgrounds)
            master_files = [f for f in file_list if 'slideMaster' in f and f.endswith('.xml')]
            print(f"      found slide master files: {master_files}")
            
            # Store raw theme XML for complete preservation
            for theme_file in theme_files:
                try:
                    with zip_file.open(theme_file) as xml_file:
                        result['theme_raw'][theme_file] = xml_file.read().decode('utf-8')
                except:
                    pass
            
            # Parse primary theme file
            for theme_file in theme_files:
                try:
                    print(f"      parsing {theme_file}...")
                    
                    with zip_file.open(theme_file) as xml_file:
                        xml_content = xml_file.read()
                        root = ET.fromstring(xml_content)
                        
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        
                        # FONTS
                        font_scheme = root.find('.//a:fontScheme', ns)
                        if font_scheme is not None:
                            major_latin = font_scheme.find('.//a:majorFont/a:latin', ns)
                            if major_latin is not None and major_latin.get('typeface'):
                                result['fonts']['title']['name'] = major_latin.get('typeface')
                                print(f"      ✓ title font: {major_latin.get('typeface')}")
                            
                            minor_latin = font_scheme.find('.//a:minorFont/a:latin', ns)
                            if minor_latin is not None and minor_latin.get('typeface'):
                                result['fonts']['body']['name'] = minor_latin.get('typeface')
                                print(f"      ✓ body font: {minor_latin.get('typeface')}")
                        
                        # COLOR SCHEME
                        color_scheme = root.find('.//a:clrScheme', ns)
                        if color_scheme is not None:
                            print(f"      ✓ found color scheme")
                            color_count = 0
                            for color_elem in color_scheme:
                                color_name = color_elem.tag.split('}')[-1]  # Remove namespace
                                # Extract color value (srgbClr or sysClr)
                                srgb = color_elem.find('.//a:srgbClr', ns)
                                sysclr = color_elem.find('.//a:sysClr', ns)
                                
                                if srgb is not None and srgb.get('val'):
                                    result['color_scheme'][color_name] = {
                                        'type': 'rgb',
                                        'value': srgb.get('val')
                                    }
                                    color_count += 1
                                elif sysclr is not None and sysclr.get('val'):
                                    result['color_scheme'][color_name] = {
                                        'type': 'system',
                                        'value': sysclr.get('val')
                                    }
                                    color_count += 1
                            
                            print(f"      ✓ extracted {color_count} theme colors")
                        
                        # FORMAT SCHEME (fill, line, effect styles)
                        format_scheme = root.find('.//a:fmtScheme', ns)
                        if format_scheme is not None:
                            print(f"      ✓ found format scheme")
                            
                            # Store fill styles
                            fill_style_list = format_scheme.find('.//a:fillStyleLst', ns)
                            if fill_style_list is not None:
                                result['format_scheme']['fills'] = len(list(fill_style_list))
                                print(f"      ✓ {len(list(fill_style_list))} fill styles")
                            
                            # Store line styles
                            line_style_list = format_scheme.find('.//a:lnStyleLst', ns)
                            if line_style_list is not None:
                                result['format_scheme']['lines'] = len(list(line_style_list))
                                print(f"      ✓ {len(list(line_style_list))} line styles")
                            
                            # Store effect styles
                            effect_style_list = format_scheme.find('.//a:effectStyleLst', ns)
                            if effect_style_list is not None:
                                result['format_scheme']['effects'] = len(list(effect_style_list))
                                print(f"      ✓ {len(list(effect_style_list))} effect styles")
                        
                        break  # Use first theme file
                
                except Exception as e:
                    print(f"      error parsing {theme_file}: {e}")
                    continue
            
            # Extract slide master backgrounds
            for master_file in master_files[:1]:  # Just get primary master
                try:
                    with zip_file.open(master_file) as xml_file:
                        xml_content = xml_file.read()
                        root = ET.fromstring(xml_content)
                        
                        ns_p = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        
                        # Check for background
                        bg = root.find('.//p:bg', ns_p)
                        if bg is not None:
                            result['backgrounds'].append({
                                'source': 'slide_master',
                                'has_background': True
                            })
                            print(f"      ✓ found slide master background")
                
                except Exception as e:
                    print(f"      error parsing {master_file}: {e}")
    
    except Exception as e:
        print(f"      ❌ CRITICAL: error extracting theme data: {e}")
        import traceback
        traceback.print_exc()
        raise ValueError(f"Failed to extract theme data from PowerPoint file: {str(e)}")
    
    # STRICT VALIDATION: Ensure we extracted real theme data
    validation_errors = []
    
    # Check fonts
    if result['fonts']['title']['name'] == 'Calibri' and result['fonts']['body']['name'] == 'Calibri':
        validation_errors.append("Failed to extract real font names (got defaults)")
    
    # Check colors
    if not result['color_scheme'] or len(result['color_scheme']) == 0:
        validation_errors.append("Failed to extract color scheme")
    
    # Check theme raw XML
    if not result['theme_raw'] or len(result['theme_raw']) == 0:
        validation_errors.append("Failed to extract theme XML data")
    
    if validation_errors:
        error_msg = "Theme extraction incomplete:\n  - " + "\n  - ".join(validation_errors)
        error_msg += "\n\nPlease ensure your PowerPoint template has a properly defined theme with fonts and colors."
        raise ValueError(error_msg)
    
    # Summary
    theme_summary = []
    if result['color_scheme']:
        theme_summary.append(f"{len(result['color_scheme'])} colors")
    if result['format_scheme']:
        theme_summary.append("format styles")
    if result['backgrounds']:
        theme_summary.append("backgrounds")
    
    if theme_summary:
        print(f"      ✅ THEME EXTRACTION COMPLETE: {', '.join(theme_summary)}")
    
    return result


def get_default_fonts_from_master(slide, pptx_path=None):
    """
    extracts default font properties from the slide master/theme.
    returns both title and body fonts.
    """
    result = {
        'title': {'name': 'Calibri', 'size': 44},
        'body': {'name': 'Calibri', 'size': 18}
    }
    
    # try method 0: direct ZIP extraction (most reliable)
    if pptx_path:
        theme_data = extract_complete_theme_from_zip(pptx_path)
        fonts = theme_data.get('fonts', {})
        if fonts.get('title', {}).get('name') != 'Calibri' or fonts.get('body', {}).get('name') != 'Calibri':
            print(f"      ✓ successfully extracted fonts from theme")
            return fonts
        else:
            print(f"      theme extraction didn't find fonts, trying other methods...")
    
    try:
        slide_layout = slide.slide_layout
        slide_master = slide_layout.slide_master
        
        print(f"      attempting to extract fonts from master slide...")
        
        # try method 1: access theme XML directly via part
        try:
            print(f"      checking slide_master.part...")
            print(f"      has 'part' attr: {hasattr(slide_master, 'part')}")
            
            if hasattr(slide_master, 'part'):
                print(f"      slide_master.part exists: {slide_master.part}")
                print(f"      has 'theme' attr: {hasattr(slide_master.part, 'theme')}")
                
                if hasattr(slide_master.part, 'theme'):
                    theme_part = slide_master.part.theme
                    print(f"      found theme part: {theme_part}")
                    
                    if hasattr(theme_part, 'element'):
                        theme_element = theme_part.element
                        print(f"      found theme element: {type(theme_element)}")
                        
                        # parse font scheme from XML
                        # namespace for theme elements
                        nsmap = {
                            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
                        }
                        
                        # find fontScheme element
                        font_scheme = theme_element.find('.//a:fontScheme', nsmap)
                        if font_scheme is not None:
                            print(f"      found fontScheme in XML")
                            
                            # extract major (title) font
                            major_font = font_scheme.find('.//a:majorFont/a:latin', nsmap)
                            if major_font is not None and major_font.get('typeface'):
                                result['title']['name'] = major_font.get('typeface')
                                print(f"      ✓ extracted theme title font from XML: {major_font.get('typeface')}")
                            
                            # extract minor (body) font
                            minor_font = font_scheme.find('.//a:minorFont/a:latin', nsmap)
                            if minor_font is not None and minor_font.get('typeface'):
                                result['body']['name'] = minor_font.get('typeface')
                                print(f"      ✓ extracted theme body font from XML: {minor_font.get('typeface')}")
                        else:
                            print(f"      fontScheme not found in XML")
                    else:
                        print(f"      theme_part has no 'element' attribute")
                else:
                    print(f"      slide_master.part has no 'theme' attribute")
            else:
                print(f"      slide_master has no 'part' attribute")
        except Exception as e:
            print(f"      error parsing theme XML: {e}")
            import traceback
            traceback.print_exc()
        
        # try method 2: legacy API approach (kept as fallback)
        if result['title']['name'] == 'Calibri' and hasattr(slide_master, 'theme') and slide_master.theme:
            theme = slide_master.theme
            print(f"      trying legacy API approach...")
            if hasattr(theme, 'font_scheme') and theme.font_scheme:
                font_scheme = theme.font_scheme
                
                # get major (heading/title) font
                if hasattr(font_scheme, 'major_font') and font_scheme.major_font:
                    major_font = font_scheme.major_font
                    if hasattr(major_font, 'latin') and major_font.latin:
                        result['title']['name'] = major_font.latin
                        print(f"      ✓ extracted theme title font via API: {major_font.latin}")
                
                # get minor (body) font
                if hasattr(font_scheme, 'minor_font') and font_scheme.minor_font:
                    minor_font = font_scheme.minor_font
                    if hasattr(minor_font, 'latin') and minor_font.latin:
                        result['body']['name'] = minor_font.latin
                        print(f"      ✓ extracted theme body font via API: {minor_font.latin}")
        
        # try method 3: extract from actual text in master slide placeholders
        if result['title']['name'] == 'Calibri' or result['body']['name'] == 'Calibri':
            print(f"      fonts still not found, trying to extract from master slide shapes...")
            for shape in slide_master.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    if shape.text_frame.paragraphs:
                        para = shape.text_frame.paragraphs[0]
                        if not para.runs:
                            para.add_run().text = "X"
                        if para.runs:
                            font = para.runs[0].font
                            font_name = font.name
                            font_size = font.size.pt if font.size else None
                            
                            if font_name:
                                # assume top shapes are title, bottom are body
                                if shape.top < (5143500 * 0.2):
                                    result['title']['name'] = font_name
                                    if font_size:
                                        result['title']['size'] = font_size
                                    print(f"      ✓ extracted title font from master shape: {font_name} {font_size}pt")
                                else:
                                    result['body']['name'] = font_name
                                    if font_size:
                                        result['body']['size'] = font_size
                                    print(f"      ✓ extracted body font from master shape: {font_name} {font_size}pt")
                                break
        
    except Exception as e:
        print(f"      ❌ CRITICAL: failed to extract theme fonts: {e}")
        import traceback
        traceback.print_exc()
        raise ValueError(f"Failed to extract font information from presentation theme: {str(e)}")
    
    # STRICT VALIDATION: Fail if we still have default fonts
    if result['title']['name'] == 'Calibri' and result['body']['name'] == 'Calibri':
        raise ValueError(
            "Failed to extract real font information from presentation. "
            "The template does not have properly defined theme fonts. "
            "Please ensure your PowerPoint template has a valid theme with defined fonts."
        )
    
    return result


def classify_layout_as_special(layout_def):
    """
    automatically classify a layout as special based on complexity indicators:
    - many text placeholders (> 3)
    - many image placeholders (> 2)
    - very small text boxes (< 15% of slide width or height)
    - very small image placeholders (< 15% of slide width or height)
    """
    placeholders = layout_def.get('placeholders', [])
    
    text_placeholders = [p for p in placeholders if p['type'] == 'text']
    image_placeholders = [p for p in placeholders if p['type'] == 'image']
    
    # criterion 1: many text boxes
    if len(text_placeholders) > 3:
        print(f"    → marked as SPECIAL: {len(text_placeholders)} text placeholders (> 3)")
        return True
    
    # criterion 2: many images
    if len(image_placeholders) > 2:
        print(f"    → marked as SPECIAL: {len(image_placeholders)} image placeholders (> 2)")
        return True
    
    # criterion 3 & 4: check for very small placeholders
    # standard slide dimensions in EMUs (16:9)
    standard_width = 9144000  # 10 inches
    standard_height = 5143500  # 5.625 inches
    
    small_threshold_width = standard_width * 0.15  # 15% of slide width
    small_threshold_height = standard_height * 0.15  # 15% of slide height
    
    for ph in placeholders:
        position = ph.get('position', {})
        width = position.get('width', 0)
        height = position.get('height', 0)
        
        if width > 0 and height > 0:
            width_pct = (width / standard_width) * 100
            height_pct = (height / standard_height) * 100
            if width < small_threshold_width or height < small_threshold_height:
                print(f"    → marked as SPECIAL: placeholder '{ph['name']}' is very small ({width_pct:.1f}% x {height_pct:.1f}%)")
                return True
    
    return False


def extract_slide_background(slide):
    background_info = {
        'type': 'none',
        'fill': None,
        'has_background_graphics': False
    }
    
    try:
        # Check if slide has a background
        if hasattr(slide, 'background') and slide.background:
            bg = slide.background
            background_info['has_background'] = True
            
            # Extract fill information
            if hasattr(bg, 'fill'):
                fill = bg.fill
                if hasattr(fill, 'type'):
                    background_info['type'] = str(fill.type)
                    
                    # Solid fill
                    if hasattr(fill, 'fore_color') and fill.fore_color:
                        try:
                            if hasattr(fill.fore_color, 'rgb'):
                                rgb = fill.fore_color.rgb
                                background_info['fill'] = {
                                    'type': 'solid',
                                    'color': {
                                        'rgb': f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                                    }
                                }
                        except:
                            pass
        
        # Check for background graphics (from layout or master)
        if hasattr(slide, 'follow_master_background'):
            background_info['follow_master_background'] = slide.follow_master_background
        
    except Exception as e:
        print(f"    warning: failed to extract background: {e}")
    
    return background_info


def extract_slide_as_layout(slide, layout_index, pptx_path=None):
    """
    extracts layout from a slide, treating all text boxes and images as content areas.
    """
    slide_layout = slide.slide_layout
    layout_name = slide_layout.name
    
    default_fonts = get_default_fonts_from_master(slide, pptx_path)
    print(f"  default fonts from master - title: {default_fonts['title']}, body: {default_fonts['body']}")
    
    # try to find better name if it's just "blank"
    if layout_name.lower() == 'blank':
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                    layout_name = shape.text_frame.text.strip()[:30]
                    break
            except:
                pass
    
    # Extract slide background
    background_info = extract_slide_background(slide)
    
    layout_def = {
        'name': f"{layout_name} (slide {layout_index + 1})",
        'original_layout_name': layout_name,
        'layout_index': layout_index,
        'slide_number': layout_index + 1,
        'placeholders': [],
        'shapes': [],
        'background': background_info,
        'category': None,
        'category_confidence': None,
        'category_rationale': None
    }
    
    internal_counter = 0
    
    for shape in slide.shapes:
        try:
            # determine if this is likely a TITLE/SUBTITLE placeholder
            #
            # IMPORTANT:
            # Many templates have placeholder runs with font.size == None, so we fall back to theme defaults.
            # If we choose the wrong default (body vs title), the AI will massively mis-estimate capacity:
            # it will think a title box is 18pt and put paragraphs in it, but PowerPoint renders it at 44pt.
            #
            # So we must classify title placeholders using placeholder type/name (NOT only y-position).
            is_likely_title = False
            try:
                if getattr(shape, "is_placeholder", False) and hasattr(shape, "placeholder_format"):
                    placeholder_type = shape.placeholder_format.type
                    # PP_PLACEHOLDER.* varies by python-pptx version; compare by name too
                    pt_str = str(placeholder_type).upper()
                    if "TITLE" in pt_str or "CENTER_TITLE" in pt_str or "SUBTITLE" in pt_str:
                        is_likely_title = True
            except Exception:
                pass

            # fallback: shape naming conventions used by PowerPoint
            try:
                name_u = (shape.name or "").upper()
                if "TITLE" in name_u or "SUBTITLE" in name_u:
                    is_likely_title = True
            except Exception:
                pass

            # final fallback: y-position heuristic (less reliable across templates/sizes)
            if not is_likely_title:
                try:
                    slide_height = 5143500  # standard 16:9 height in EMUs (best-effort)
                    is_likely_title = shape.top < (slide_height * 0.25)  # slightly less strict than 20%
                except Exception:
                    is_likely_title = False
            
            default_font = default_fonts['title'] if is_likely_title else default_fonts['body']
            print(f"      shape '{shape.name}' at y={shape.top}, using {'title' if is_likely_title else 'body'} font defaults")
            shape_props = extract_shape_complete_properties(shape, default_font)
            
            # all text frames and images are content areas
            if shape_props.get('is_placeholder', False):
                content_type = shape_props['content_type']
                layout_def['placeholders'].append({
                    'idx': internal_counter,
                    'real_pptx_idx': shape_props['placeholder_idx'],
                    'type': content_type,
                    'name': shape.name,
                    'position': shape_props['position'],
                    'properties': shape_props
                })
                print(f"      → added {content_type.upper()} placeholder idx={internal_counter}: {shape.name}")
                internal_counter += 1
            else:
                # Static shape (not a placeholder for user content)
                layout_def['shapes'].append(shape_props)
                shape_type_str = shape_props.get('shape_type', 'UNKNOWN')
                if shape_props.get('is_design_image'):
                    print(f"      → added DESIGN IMAGE (static): {shape.name} [{shape_type_str}]")
                else:
                    print(f"      → added STATIC SHAPE: {shape.name} [{shape_type_str}]")
        
        except Exception as e:
            print(f"  warning: failed to extract shape '{shape.name}': {e}")
            import traceback
            traceback.print_exc()
            continue
    
    # summary of what was extracted
    text_count = len([p for p in layout_def['placeholders'] if p['type'] == 'text'])
    image_count = len([p for p in layout_def['placeholders'] if p['type'] == 'image'])
    static_count = len(layout_def['shapes'])
    print(f"    ✓ layout complete: {text_count} text + {image_count} image placeholders, {static_count} static shapes")
    
    return layout_def


def extract_all_slides_as_layouts(presentation, pptx_path=None):
    """
    extracts layouts from all slides in presentation.
    also extracts complete theme data for preservation.
    """
    # Extract complete theme data ONCE at the beginning
    theme_data = None
    if pptx_path:
        print(f"\n" + "="*80)
        print(f"EXTRACTING COMPLETE THEME DATA")
        print(f"="*80)
        theme_data = extract_complete_theme_from_zip(pptx_path)
        print(f"="*80 + "\n")
    
    layouts = []
    print(f"\nextracting {len(presentation.slides)} slides as layout templates...")
    
    for idx, slide in enumerate(presentation.slides):
        print(f"  extracting slide {idx + 1} ({slide.slide_layout.name})...")
        layout = extract_slide_as_layout(slide, idx, pptx_path)
        
        print(f"    found {len(layout['placeholders'])} content areas")
        print(f"    found {len(layout['shapes'])} static shapes")
        
        layouts.append(layout)
    
    total_text_placeholders = sum(len([p for p in l['placeholders'] if p['type'] == 'text']) for l in layouts)
    total_image_placeholders = sum(len([p for p in l['placeholders'] if p['type'] == 'image']) for l in layouts)
    layouts_with_images = sum(1 for l in layouts if any(p['type'] == 'image' for p in l['placeholders']))
    
    # count all static shapes
    total_static_shapes = sum(len(l['shapes']) for l in layouts)
    total_design_images = sum(len([s for s in l['shapes'] if s.get('is_design_image')]) for l in layouts)
    layouts_with_design_images = sum(1 for l in layouts if any(s.get('is_design_image') for s in l['shapes']))
    
    # count shape types
    shape_types = {}
    for l in layouts:
        for s in l['shapes']:
            shape_type = s.get('shape_type', 'UNKNOWN')
            shape_types[shape_type] = shape_types.get(shape_type, 0) + 1
    
    print(f"\n{'='*80}")
    print(f"EXTRACTION COMPLETE")
    print(f"{'='*80}")
    print(f"Total layouts extracted: {len(layouts)}")
    print(f"Layouts with IMAGE support: {layouts_with_images}")
    print(f"Layouts with TEXT only: {len(layouts) - layouts_with_images}")
    print(f"Total text placeholders: {total_text_placeholders}")
    print(f"Total image placeholders: {total_image_placeholders}")
    print(f"Total static shapes: {total_static_shapes}")
    print(f"  - Design images: {total_design_images}")
    print(f"  - Other shapes: {total_static_shapes - total_design_images}")
    
    if shape_types:
        print(f"\nStatic shape types found:")
        for shape_type, count in sorted(shape_types.items(), key=lambda x: -x[1]):
            print(f"  - {shape_type}: {count}")
    
    print(f"{'='*80}\n")
    
    if layouts_with_images == 0:
        print("⚠️  WARNING: No image placeholders detected in any layout!")
        print("   Make sure your template has slides with image placeholders (not just existing images).")
        print("   Image placeholders should be empty placeholder boxes ready to receive images.\n")
    
    if total_static_shapes > 0:
        print(f"✓ Found {total_static_shapes} static shapes that will be recreated in generated slides")
        print(f"  (including {total_design_images} design images)\n")
    
    # Return both layouts and complete theme data
    return {
        'layouts': layouts,
        'theme_data': theme_data
    }
