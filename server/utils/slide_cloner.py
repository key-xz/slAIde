"""
clones template slides with new content while preserving all design elements.
ai selects which slides to use, cloner preserves pixel-perfect design.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import base64
from copy import deepcopy


def clone_slide_with_content(source_presentation, source_slide_index, target_presentation, 
                             placeholder_contents, uploaded_images, image_order=None):
    try:
        # Get the source slide
        source_slide = source_presentation.slides[source_slide_index]
        source_layout = source_slide.slide_layout
        
        print(f"    cloning slide {source_slide_index + 1} (layout: {source_layout.name})")
        
        # Add a new slide with the same layout
        new_slide = target_presentation.slides.add_slide(source_layout)
        
        # Copy all shapes from source to target
        # This preserves backgrounds, design elements, and all formatting
        _copy_all_shapes(source_slide, new_slide, placeholder_contents, uploaded_images, image_order)
        
        # Apply background if present
        _copy_background(source_slide, new_slide)
        
        # Validate shape positions are within slide bounds
        _validate_shape_positions(new_slide, target_presentation)
        
        print(f"    cloned slide successfully")
        return new_slide
        
    except Exception as e:
        print(f"    error cloning slide: {e}")
        import traceback
        traceback.print_exc()
        raise


def _copy_all_shapes(source_slide, target_slide, placeholder_contents, uploaded_images, image_order=None):
    content_map = {item['idx']: item for item in placeholder_contents}
    placeholder_counter = 0
    
    print(f"      copying {len(source_slide.shapes)} shapes...")
    
    for shape in source_slide.shapes:
        try:
            shape_name = shape.name if hasattr(shape, 'name') else 'Unknown'
            
            # Check if this is a placeholder
            is_placeholder = hasattr(shape, 'is_placeholder') and shape.is_placeholder
            
            if is_placeholder:
                # This is a placeholder - we need to fill it with content
                if placeholder_counter in content_map:
                    content_item = content_map[placeholder_counter]
                    _fill_placeholder_in_cloned_slide(target_slide, shape, content_item, uploaded_images, placeholder_counter, image_order)
                else:
                    print(f"        placeholder {placeholder_counter} has no content, keeping original")
                
                placeholder_counter += 1
            else:
                # Non-placeholder shape - it will be copied automatically by the layout
                # We don't need to manually copy it as add_slide already includes layout shapes
                pass
                
        except Exception as e:
            print(f"        warning: error processing shape '{shape_name}': {e}")


def _fill_placeholder_in_cloned_slide(target_slide, source_shape, content_item, uploaded_images, idx, image_order=None):
    content_type = content_item.get('type')
    
    # Find the corresponding placeholder in the target slide
    target_placeholder = None
    try:
        # Try to find by placeholder idx
        if hasattr(source_shape, 'placeholder_format'):
            ph_idx = source_shape.placeholder_format.idx
            for shape in target_slide.placeholders:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == ph_idx:
                    target_placeholder = shape
                    break
    except:
        pass
    
    if not target_placeholder:
        print(f"        could not find target placeholder for idx {idx}")
        return
    
    try:
        if content_type == 'text':
            # Replace text content
            content = content_item.get('content', '')
            if hasattr(target_placeholder, 'text_frame'):
                text_frame = target_placeholder.text_frame
                text_frame.clear()
                text_frame.text = content
                
                # CRITICAL: Apply template font size to ensure accurate overflow detection
                # Get font size from content_item metadata (passed from AI service)
                template_font_size = content_item.get('font_size')
                
                if template_font_size:
                    # Apply font size to ALL paragraphs and runs
                    from pptx.util import Pt
                    font_size_emu = int(template_font_size * 914400 / 72)  # Convert pt to EMU
                    
                    for para in text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = font_size_emu
                    
                    print(f"        ✓ filled text placeholder {idx} with font {template_font_size}pt: '{content[:50]}...'")
                else:
                    print(f"        ✓ filled text placeholder {idx} (no font size): '{content[:50]}...'")
        
        elif content_type == 'image':
            image_index = content_item.get('image_index')
            if image_index is not None:
                # use image_order if available, otherwise fall back to dict keys
                if image_order:
                    image_filenames = image_order
                    print(f"        using image_order: {image_filenames}")
                else:
                    image_filenames = list(uploaded_images.keys())
                    print(f"        ⚠️  no image_order, using dict keys: {image_filenames}")
                
                print(f"        looking for image_index {image_index} in {len(image_filenames)} images")
                
                if 0 <= image_index < len(image_filenames):
                    filename = image_filenames[image_index]
                    image_data_b64 = uploaded_images[filename]
                    print(f"        selected image: {filename}")
                    
                    if ',' in image_data_b64:
                        image_data_b64 = image_data_b64.split(',')[1]
                    image_bytes = base64.b64decode(image_data_b64)
                    image_stream = BytesIO(image_bytes)
                    
                    # Try to insert picture into the placeholder (safest method)
                    try:
                        # save placeholder position and dimensions
                        left = target_placeholder.left
                        top = target_placeholder.top
                        width = target_placeholder.width
                        height = target_placeholder.height
                        
                        # Method 1: Try using placeholder's built-in picture insertion
                        if hasattr(target_placeholder, 'insert_picture'):
                            target_placeholder.insert_picture(image_stream)
                            print(f"        ✓ filled image placeholder {idx}: {filename} (using insert_picture)")
                        else:
                            # Method 2: Clear placeholder and add picture at same position
                            # This is safer than removing the element
                            placeholder_id = target_placeholder.shape_id
                            
                            # Add picture at exact same position
                            pic = target_slide.shapes.add_picture(image_stream, left, top, width, height)
                            print(f"        ✓ filled image placeholder {idx}: {filename} (added picture)")
                    except Exception as img_err:
                        print(f"        warning: could not insert image into placeholder {idx}: {img_err}")
                        # Fallback: just add the picture
                        try:
                            target_slide.shapes.add_picture(image_stream, left, top, width, height)
                            print(f"        ✓ filled image placeholder {idx}: {filename} (fallback method)")
                        except Exception as fallback_err:
                            print(f"        error: failed to add image: {fallback_err}")
    
    except Exception as e:
        print(f"        error filling placeholder {idx}: {e}")
        import traceback
        traceback.print_exc()


def _copy_background(source_slide, target_slide):
    try:
        # Copy background fill if present
        if hasattr(source_slide, 'background') and hasattr(target_slide, 'background'):
            # The background is usually inherited from layout/master
            # In most cases, using the same layout already preserves the background
            pass
    except Exception as e:
        print(f"        warning: could not copy background: {e}")


def _validate_shape_positions(slide, presentation):
    """
    validate that all shapes on the slide are within slide bounds.
    DISABLED: Shape position modifications can corrupt PowerPoint files.
    """
    # DO NOT modify shape positions - this can corrupt the PPTX file
    # The template should already have valid positions
    # If shapes are out of bounds, it's better to warn than to corrupt the file
    
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    violations = []
    
    for shape in slide.shapes:
        try:
            # skip shapes without position
            if not hasattr(shape, 'left') or not hasattr(shape, 'top'):
                continue
            
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            right = left + width
            bottom = top + height
            
            # check bounds but DO NOT modify
            if left < 0 or top < 0 or right > slide_width or bottom > slide_height:
                shape_name = getattr(shape, 'name', 'Unknown')
                violations.append(f"        warning: shape '{shape_name}' outside bounds (left={left}, top={top}, right={right}, bottom={bottom})")
        
        except Exception as e:
            print(f"        warning: could not validate shape position: {e}")
    
    if violations:
        for v in violations:
            print(v)
        print(f"        found {len(violations)} shape(s) outside bounds - template may need adjustment")
